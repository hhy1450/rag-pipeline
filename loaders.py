"""
Multi-format document loaders.

支持格式: PDF / Word(.docx) / PPT(.pptx) / 图片(OCR) / Markdown / TXT

用法:
    loader = get_loader("file.pdf")
    text = loader.load()
"""
from pathlib import Path


# ---- PDF ----
class PDFLoader:
    """PyMuPDF 提取 PDF 纯文本。"""
    def __init__(self, path: str | Path):
        self.path = Path(path)

    def load(self) -> str:
        import fitz
        doc = fitz.open(str(self.path))
        parts = []
        for page in doc:
            text = page.get_text()
            if text.strip():
                parts.append(text)
        doc.close()
        return "\n\n".join(parts)


# ---- Word (.docx) ----
class DocxLoader:
    """python-docx 提取 Word 文档纯文本。"""
    def __init__(self, path: str | Path):
        self.path = Path(path)

    def load(self) -> str:
        from docx import Document
        doc = Document(str(self.path))
        parts = []
        for para in doc.paragraphs:
            if para.text.strip():
                parts.append(para.text)
        # Also extract tables
        for table in doc.tables:
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if cells:
                    parts.append(" | ".join(cells))
        return "\n\n".join(parts)


# ---- PPT (.pptx) ----
class PPTLoader:
    """python-pptx 提取 PPT 文本（含标题+正文+表格）。"""
    def __init__(self, path: str | Path):
        self.path = Path(path)

    def load(self) -> str:
        from pptx import Presentation
        prs = Presentation(str(self.path))
        parts = []
        for i, slide in enumerate(prs.slides):
            slide_parts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text_frame.text.strip()
                    if text:
                        slide_parts.append(text)
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                        if cells:
                            slide_parts.append(" | ".join(cells))
            if slide_parts:
                parts.append(f"[Slide {i+1}]\n" + "\n".join(slide_parts))
        return "\n\n".join(parts)


# ---- Image (OCR) ----
class ImageLoader:
    """Pillow + RapidOCR 提取图片中的文字。

    依赖: pip install rapidocr-onnxruntime
    """
    def __init__(self, path: str | Path):
        self.path = Path(path)

    def load(self) -> str:
        try:
            from rapidocr_onnxruntime import RapidOCR
            engine = RapidOCR()
            result, _ = engine(str(self.path))
            if result:
                lines = []
                for (_, text, confidence) in result:
                    if text.strip() and confidence > 0.5:
                        lines.append(text.strip())
                return "\n".join(lines)
            return ""
        except ImportError:
            # Fallback: return filename as placeholder
            return f"[图片文件: {self.path.name} (安装 rapidocr-onnxruntime 以启用 OCR)]"
        except Exception as e:
            return f"[OCR 失败: {e}]"


# ---- Markdown / TXT ----
class TextLoader:
    """加载纯文本文件 (Markdown, TXT, 代码等)。"""
    def __init__(self, path: str | Path):
        self.path = Path(path)

    def load(self) -> str:
        return self.path.read_text(encoding="utf-8", errors="replace")


# ---- Loader factory ----
LOADER_MAP = {
    ".pdf": PDFLoader,
    ".docx": DocxLoader,
    ".pptx": PPTLoader,
    ".ppt": PPTLoader,
    ".png": ImageLoader,
    ".jpg": ImageLoader,
    ".jpeg": ImageLoader,
    ".bmp": ImageLoader,
    ".md": TextLoader,
    ".txt": TextLoader,
    ".csv": TextLoader,
}


def get_loader(path: str | Path) -> object:
    """根据文件后缀自动选择 Loader。"""
    suffix = Path(path).suffix.lower()
    loader_cls = LOADER_MAP.get(suffix)
    if loader_cls is None:
        raise ValueError(f"Unsupported format: {suffix}")
    return loader_cls(path)


def load_directory(directory: str | Path) -> list[dict]:
    """加载目录下所有支持格式的文件。

    Returns:
        [{"filename": "xxx.pdf", "text": "..."}, ...]
    """
    docs = []
    dir_path = Path(directory)
    for ext in LOADER_MAP:
        for file_path in dir_path.glob(f"*{ext}"):
            try:
                loader = get_loader(file_path)
                docs.append({
                    "filename": file_path.name,
                    "text": loader.load(),
                })
            except Exception as e:
                docs.append({
                    "filename": file_path.name,
                    "text": f"[加载失败: {e}]",
                })
    return docs
